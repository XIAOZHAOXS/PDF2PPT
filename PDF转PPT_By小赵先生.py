"""
Author:赵文瑄
Date:2023.12.10
Power By Pycharm
"""
import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
import tempfile
from time import sleep

def convert_pdf_to_ppt(pdf_path, ppt_path, dpi=300):
    # 打开PDF文件
    pdf = fitz.open(pdf_path)

    # 创建一个PPT演示文稿
    ppt = Presentation()

    # 将每一页PDF转换为PPT中的一页幻灯片
    for page_number in range(len(pdf)):
        page = pdf.load_page(page_number)
        # 提高图像的分辨率
        pix = page.get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72))
        img_data = pix.tobytes("png")

        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # 使用空白幻灯片布局
        left = top = Inches(0)

        # 使用临时文件
        with BytesIO(img_data) as img_stream:
            slide.shapes.add_picture(img_stream, left, top, width=ppt.slide_width)

    # 保存PPT
    ppt.save(ppt_path)
    pdf.close()

# 用户可以通过命令行参数指定输出路径或使用默认路径
output_dir = "PPT"  # 可以根据需要更改
if not os.path.exists(output_dir):
    os.makedirs(output_dir)

current_dir = os.getcwd()
ppt_dir = os.path.join(current_dir, output_dir)

# 将每个PDF转换为PPT
for pdf_file in os.listdir(current_dir):
    if pdf_file.endswith('.pdf'):
        pdf_path = os.path.join(current_dir, pdf_file)
        ppt_file_name = os.path.splitext(pdf_file)[0] + '.pptx'
        ppt_path = os.path.join(ppt_dir, ppt_file_name)
        convert_pdf_to_ppt(pdf_path, ppt_path, dpi=300)
        print(f"已将 '{pdf_file}' 转换为 '{ppt_file_name}'")

print("所有PDF文件已成功转换为PPT演示文稿。")
print("\033[31m" + "Power By 小赵先生" + "\033[0m")
sleep(5)